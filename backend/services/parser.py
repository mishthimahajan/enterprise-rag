# import pdfplumber

# from docx import Document
# from pptx import Presentation


# def parse_pdf(path: str) -> str:
#     """
#     Extract text from PDF using pdfplumber
#     """
#     text = ""

#     with pdfplumber.open(path) as pdf:

#         for page in pdf.pages:

#             page_text = page.extract_text()

#             if page_text:
#                 text += page_text + "\n"

#     return text.strip()


# def parse_docx(path: str) -> str:
#     """
#     Extract text from Word document
#     """
#     doc = Document(path)

#     text = ""

#     for para in doc.paragraphs:

#         if para.text.strip():
#             text += para.text + "\n"

#     return text.strip()


# def parse_ppt(path: str) -> str:
#     """
#     Extract text from PowerPoint
#     """
#     prs = Presentation(path)

#     text = ""

#     for slide in prs.slides:

#         for shape in slide.shapes:

#             if hasattr(shape, "text"):

#                 if shape.text.strip():
#                     text += shape.text + "\n"

#     return text.strip()


import re

from pypdf import PdfReader
from docx import Document
from pptx import Presentation


def clean_text(text: str):

    text = re.sub(
        r"\s+",
        " ",
        text
    )

    return text.strip()


# PDF Parser
def parse_pdf(path: str):

    text = ""

    reader = PdfReader(path)

    for page in reader.pages:

        page_text = page.extract_text()

        if page_text:
            text += page_text + "\n"

    return clean_text(text)


# DOCX Parser
def parse_docx(path: str):

    doc = Document(path)

    text = "\n".join(
        para.text
        for para in doc.paragraphs
    )

    return clean_text(text)


# PPT Parser
def parse_ppt(path: str):

    prs = Presentation(path)

    text = ""

    for slide in prs.slides:

        for shape in slide.shapes:

            if hasattr(shape, "text"):

                text += shape.text + "\n"


    return clean_text(text)